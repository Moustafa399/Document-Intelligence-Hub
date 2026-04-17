"""PowerPoint file extractor."""
import os
from pptx import Presentation

from helpers.file_utils import create_document_folder, save_text, save_metadata, save_tables
from helpers.table_utils import format_table_as_markdown


def extract_ppt(file_path):
    """Extract text, tables, and images from PowerPoint presentation."""
    doc_id, base, text_dir, img_dir = create_document_folder(file_path)

    prs = Presentation(file_path)
    text = ""
    images = []
    tables_data = []
    counter = 1

    for slide_num, slide in enumerate(prs.slides, 1):
        text += f"\n\n=== SLIDE {slide_num} ===\n"
        
        for shape in slide.shapes:
            # Extract text from shapes
            if hasattr(shape, "text") and shape.text.strip():
                text += shape.text + "\n"
            
            # Check if shape is a table
            if shape.shape_type == 19:  # MSO_SHAPE_TYPE.TABLE = 19
                try:
                    if hasattr(shape, "table"):
                        table = shape.table
                        table_data = []
                        for row in table.rows:
                            table_data.append([cell.text.strip() for cell in row.cells])
                        
                        if table_data:
                            tables_data.append({
                                "slide": slide_num,
                                "table_index": len(tables_data) + 1,
                                "data": table_data
                            })
                            
                            text += f"\n[TABLE Slide {slide_num}]\n"
                            text += format_table_as_markdown(table_data)
                            text += "\n"
                except Exception as e:
                    print(f"⚠️ Could not extract table from slide {slide_num}: {e}")
            
            # Extract images (shape_type 13 is picture)
            if shape.shape_type == 13:
                try:
                    img = shape.image
                    path = os.path.join(img_dir, f"img_{counter}.{img.ext}")
                    with open(path, "wb") as f:
                        f.write(img.blob)
                    images.append(path)
                    counter += 1
                except Exception as e:
                    print(f"⚠️ Could not extract image from slide {slide_num}: {e}")

    # Save tables separately as JSON
    if tables_data:
        save_tables(base, tables_data)
        print(f"📊 Found {len(tables_data)} table(s) in PowerPoint")

    save_text(text_dir, text)
    save_metadata(base, {
        "source": "powerpoint", 
        "images_found": len(images),
        "tables_found": len(tables_data)
    })

    return base, images, doc_id, "powerpoint"
